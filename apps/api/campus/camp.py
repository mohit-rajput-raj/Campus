import pandas as pd
from pptx import Presentation
import pdfplumber

def load_data(file_path):
    if file_path.endswith('.csv'):
        df = pd.read_csv(file_path)
    elif file_path.endswith('.xlsx'):
        df = pd.read_excel(file_path)
    elif file_path.endswith('.json'):
        df = pd.read_json(file_path)
    elif file_path.endswith('.pptx'):
        df = extract_from_pptx(file_path)
    elif file_path.endswith('.pdf'):
        df = extract_from_pdf(file_path)
    else:
        raise ValueError("Unsupported file format")
    return df

def extract_from_pptx(file_path):
    prs = Presentation(file_path)
    data = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for row in table.rows[1:]:  # skip header
                    cells = [cell.text.strip() for cell in row.cells]
                    data.append(cells)
            elif shape.has_text_frame:
                text = shape.text_frame.text
                # optionally parse lines
    return pd.DataFrame(data, columns=["Name", "Roll", "Date", "Status"])

def extract_from_pdf(file_path):
    data = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            table = page.extract_table()
            if table:
                df = pd.DataFrame(table[1:], columns=table[0])
                data.append(df)
    return pd.concat(data, ignore_index=True)

def normalize(df):
    # convert column names and data into a unified schema
    column_map = {
        "Student Name": "name",
        "Roll No": "roll_number",
        "Date": "date",
        "Attendance": "status",
    }
    df = df.rename(columns=column_map)
    df = df[[c for c in ["name", "roll_number", "date", "status"] if c in df.columns]]
    return df.to_dict(orient="records")

# Usage:
file = "AI_csv.csv"
df = load_data(file)
common_data = normalize(df)
dd = pd.DataFrame(common_data)
dd.to_csv("newdata.csv", index=False)
print(common_data)
